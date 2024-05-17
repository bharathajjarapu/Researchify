import os
import io
import PIL
import docx
import pptx
import wikipedia
import pytesseract
import pandas as pd
import streamlit as st
from pymed import PubMed
from PyPDF2 import PdfReader
from dotenv import load_dotenv
from tavily import TavilyClient
import google.generativeai as genai
from langchain.prompts import PromptTemplate
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter

load_dotenv()
tavily = TavilyClient(api_key=os.getenv('TAVILY_API_KEY'))
pubmed = PubMed(tool="MyApp", email="bharath@email.com")
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])
model = genai.GenerativeModel('gemini-pro')

st.set_page_config(
    page_title="Researchify",
    page_icon=":book:",
)

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "messages" not in st.session_state:
    st.session_state.messages = []

with st.sidebar:
    uploaded_files = st.file_uploader("Local Documents", type=["pdf", "docx", "pptx", "csv", "xls", "xlsx", "png", "jpg", "jpeg"], accept_multiple_files=True)
    ocr = st.checkbox("Enable OCR for PDF files")

def extract_text_from_file(file_bytes, file_extension, ocr=False):
    text = ""

    if file_extension == ".pdf":
        pdf_io = io.BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_io)
        if ocr:
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                images = page.images
                for image in images:
                    image_bytes = image.data
                    image_file = io.BytesIO(image_bytes)
                    image = PIL.Image.open(image_file)
                    text += pytesseract.image_to_string(image) + "\n"
        else:
            for page in pdf_reader.pages:
                text += page.extract_text()

    elif file_extension in [".png", ".jpg", ".jpeg", ".webp"]:
        image_file = io.BytesIO(file_bytes)
        image = PIL.Image.open(image_file)
        vision_model = gen_ai.GenerativeModel('gemini-pro-vision')
        response = vision_model.generate_content(["Explain the picture and Extract Text from picture?", image])
        text = response.text

    elif file_extension == ".docx":
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_extension == ".pptx":
        presentation = pptx.Presentation(io.BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_extension == ".csv":
        data = pd.read_csv(io.BytesIO(file_bytes))
        text = data.to_string()

    return text

def split_text_into_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    text_chunks = text_splitter.split_text(text)
    return text_chunks

def create_vector_store(chunks, file_names):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    metadatas = [{"source": file_name} for file_name in file_names for _ in range(len(chunks) // len(file_names))]
    vector_store = FAISS.from_texts(chunks, embedding=embeddings, metadatas=metadatas)
    st.session_state.vector_store = vector_store

def get_response(user_question):
    relevant_docs = st.session_state.vector_store.similarity_search(user_question)
    full_response = ""

    for doc in relevant_docs:
        doc_content = doc.page_content
        full_response += doc_content + "\n\n"

    full_response += "Sources : Local Documents"

    return full_response

def search_sources(topic):
    results = []

    with st.spinner("Searching Web"):
        tavily_search_results = tavily.search(topic, search_depth="advanced")
        if tavily_search_results:
            results.extend(tavily_search_results.get('results', []))

    with st.spinner("Searching PubMed"):
        pubmed_search_results = pubmed.query(topic, max_results=10)
        if pubmed_search_results:
            for result in pubmed_search_results:
                pubmed_url = f"https://www.ncbi.nlm.nih.gov/pubmed/{result.pubmed_id}"
                results.append({
                    'url': pubmed_url,
                    'content': result.abstract
                })

    with st.spinner("Searching Wikipedia"):
        try:
            wiki_page = wikipedia.page(topic)
            results.append({
                'url': wiki_page.url,
                'content': wiki_page.content[:1000] + '...'
            })
        except wikipedia.exceptions.PageError:
            pass

    with st.spinner("Searching Local Documents"):
        if st.session_state.vector_store is not None:
            try:
                results.append({
                    'url': 'Local Documents',
                    'content': get_response(topic)
                })
            except Exception as e:
                st.write(f"Error occurred while searching local documents: {e}")
            
    return results

def generate_report(topic, search_results):
    description = "You are a Senior NYT Editor tasked with writing a NYT cover story worthy report due tomorrow, Report should be as human as it can be."
    instructions = [
        "You will be provided with a topic and search results from junior researchers.",
        "Carefully read the results and generate a final - NYT cover story worthy report.",
        "Make your report engaging, informative, and well-structured.",
        "Your report should follow the format provided below.",
        "Remember: you are writing for the New York Times, so the quality of the report is important.",
    ]
    report_format = """

# Title

## **Overview** Brief introduction of the topic
## **Importance** Why is this topic significant now?

### Section 1
- **Detail 1**
- **Detail 2**
- **Detail 3**

### Section 2
- **Detail 1**
- **Detail 2**
- **Detail 3**

### Section 3
- **Detail 1**
- **Detail 2**
- **Detail 3**

### Keywords
- all the important technical keywords.

## Conclusion
- **Implications:** What these findings mean for the future.
- **Summary of report:** Recap of the key findings from the report.

## References
- [Reference 1](Link Paper or Website to Source)
Paper or Website to Source)
- [Reference 2](Link Paper or Website to Source)
- [Reference 3](Link Paper or Website to Source)
- [Reference 4](Link Paper or Websiteto Source)
- [Reference 5](Link Paper orWebsite to Source)

"""
    prompt = f"{description}\n\nInstructions: {', '.join(instructions)}\n\nReport Format:\n{report_format}\n\nTopic: {topic}\n\nSearch Results:\n"
    for result in search_results:
        prompt += f"- {result['url']}\n{result['content']}\n\n"

    chat_completion = model.generate_content(prompt)
    return chat_completion.text

def main():
    st.title(":book: Research Assistant")
    input_topic = st.text_input("Enter a research topic")
    generate_report_btn = st.button("Generate Report")

    if uploaded_files:
        file_texts = []
        file_names = []
        for uploaded_file in uploaded_files:
            file_bytes = uploaded_file.read()
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()

            st.sidebar.write(f"Processing your {uploaded_file.name} file...")
            file_text = extract_text_from_file(file_bytes, file_extension, ocr)
            file_texts.append(file_text)
            file_names.append(uploaded_file.name)
            st.sidebar.write(f"{uploaded_file.name} File processed successfully!")

        text_chunks = sum([split_text_into_chunks(text) for text in file_texts], [])
        create_vector_store(text_chunks, file_names)

    if generate_report_btn:
        st.session_state["topic"] = input_topic
        report_topic = st.session_state["topic"]
        search_results = search_sources(report_topic)

        if uploaded_files:
            for file_name, file_text in zip(file_names, file_texts):
                search_results.append({
                    'url': file_name,
                    'content': file_text
                })

        if not search_results:
            st.write("Sorry report generation failed. Please try again.")
            return

        with st.spinner("Generating Report"):
            final_report = generate_report(report_topic, search_results)
            st.markdown(final_report)

if __name__ == "__main__":
    main()
